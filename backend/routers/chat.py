"""
채팅 API 라우터
멀티 LLM을 이용한 스트리밍 채팅 API (PostgreSQL 비동기 버전)
파일 업로드/다운로드 지원 + 파일 내용 추출
"""
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form, Request
from fastapi.responses import StreamingResponse, FileResponse
from pydantic import BaseModel
from typing import Optional, List
import uuid
import json
import os
import shutil
from datetime import datetime
import base64

# 문서 파싱 라이브러리
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from llm_factory import get_llm
from routers.auth import get_current_user
from database import get_db
from sqlalchemy import text
from typing import Dict
import logging
import asyncio

# 개인화 메모리 서비스
from services.memory_service import get_memory_service

# Agent 서비스
from services.agent_service import get_agent_service

# 개인화 성능 향상 모듈 (컨설팅 반영)
from services.personalization_enhancer import (
    get_personalization_enhancer,
    QueryIntent,
    SEMANTIC_ROUTER_ENABLED,
    SELF_RAG_ENABLED,
    HIERARCHICAL_MEMORY_ENABLED,
    LLM_JUDGE_ENABLED
)

# MH_RAG 서비스 (문서 검색용)
from services.mh_rag_service import get_mh_rag_service, MH_RAG_ENABLED

# Long-term 프로파일 서비스
from services.profile_service import get_profile_service

# Long-term 메모리 설정
LONGTERM_MEMORY_ENABLED = os.getenv("LONGTERM_MEMORY_ENABLED", "true").lower() == "true"
PROFILE_AUTO_EXTRACT = os.getenv("PROFILE_AUTO_EXTRACT", "true").lower() == "true"

logger = logging.getLogger(__name__)


# ============================================================================
# 데이터베이스 함수들
# ============================================================================

async def create_session(session_id: str, user_id: int, title: str = None, model: str = "gpt4o") -> str:
    """새 대화 세션 생성"""
    async with get_db() as session:
        await session.execute(
            text("""
                INSERT INTO sessions (id, user_id, title, model)
                VALUES (:id, :user_id, :title, :model)
            """),
            {"id": session_id, "user_id": user_id, "title": title, "model": model}
        )
        return session_id


async def get_session(session_id: str) -> Dict:
    """세션 조회"""
    async with get_db() as session:
        result = await session.execute(
            text("SELECT * FROM sessions WHERE id = :id"),
            {"id": session_id}
        )
        row = result.mappings().first()
        return dict(row) if row else None


async def get_user_sessions(user_id: int) -> list:
    """사용자의 모든 세션 조회"""
    async with get_db() as session:
        result = await session.execute(
            text("""
                SELECT id, title, model, created_at, updated_at
                FROM sessions
                WHERE user_id = :user_id
                ORDER BY updated_at DESC
            """),
            {"user_id": user_id}
        )
        return [dict(row) for row in result.mappings().all()]


async def update_session_title(session_id: str, title: str):
    """세션 제목 업데이트"""
    async with get_db() as session:
        await session.execute(
            text("UPDATE sessions SET title = :title, updated_at = NOW() WHERE id = :id"),
            {"title": title, "id": session_id}
        )


async def delete_session(session_id: str):
    """세션 삭제"""
    async with get_db() as session:
        # 관련 메시지 먼저 삭제
        await session.execute(
            text("DELETE FROM messages WHERE session_id = :id"),
            {"id": session_id}
        )
        # 세션 삭제
        await session.execute(
            text("DELETE FROM sessions WHERE id = :id"),
            {"id": session_id}
        )


async def add_message(session_id: str, role: str, content: str) -> int:
    """메시지 추가"""
    async with get_db() as session:
        result = await session.execute(
            text("""
                INSERT INTO messages (session_id, role, content)
                VALUES (:session_id, :role, :content)
                RETURNING id
            """),
            {"session_id": session_id, "role": role, "content": content}
        )
        # 세션 업데이트 시간 갱신
        await session.execute(
            text("UPDATE sessions SET updated_at = NOW() WHERE id = :id"),
            {"id": session_id}
        )
        return result.scalar()


async def get_session_messages(session_id: str) -> list:
    """세션의 모든 메시지 조회"""
    async with get_db() as session:
        result = await session.execute(
            text("""
                SELECT id, role, content, created_at
                FROM messages
                WHERE session_id = :session_id
                ORDER BY created_at ASC
            """),
            {"session_id": session_id}
        )
        return [dict(row) for row in result.mappings().all()]


async def add_attachment(session_id: str, user_id: int, filename: str,
                         original_filename: str, file_type: str, file_size: int, file_path: str) -> int:
    """첨부파일 정보 저장"""
    async with get_db() as session:
        result = await session.execute(
            text("""
                INSERT INTO attachments (session_id, user_id, filename, original_filename, file_type, file_size, file_path)
                VALUES (:session_id, :user_id, :filename, :original_filename, :file_type, :file_size, :file_path)
                RETURNING id
            """),
            {
                "session_id": session_id,
                "user_id": user_id,
                "filename": filename,
                "original_filename": original_filename,
                "file_type": file_type,
                "file_size": file_size,
                "file_path": file_path
            }
        )
        return result.scalar()


async def get_attachment(attachment_id: int) -> Dict:
    """첨부파일 정보 조회"""
    async with get_db() as session:
        result = await session.execute(
            text("SELECT * FROM attachments WHERE id = :id"),
            {"id": attachment_id}
        )
        row = result.mappings().first()
        return dict(row) if row else None


async def get_session_attachments(session_id: str) -> list:
    """세션의 모든 첨부파일 조회"""
    async with get_db() as session:
        result = await session.execute(
            text("""
                SELECT id, filename, original_filename, file_type, file_size, created_at
                FROM attachments
                WHERE session_id = :session_id
                ORDER BY created_at DESC
            """),
            {"session_id": session_id}
        )
        return [dict(row) for row in result.mappings().all()]

# 파일 업로드 경로
UPLOAD_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

# 허용되는 파일 타입 (확장된 지원)
ALLOWED_EXTENSIONS = {
    # 이미지 파일
    'image': ['.jpg', '.jpeg', '.png', '.gif', '.webp'],
    # 문서 파일
    'document': ['.pdf', '.docx', '.pptx', '.hwp', '.hwpx'],
    # 표 및 데이터
    'spreadsheet': ['.csv', '.xls', '.xlsx', '.tsv'],
    # 프로그래밍 언어
    'programming': [
        '.py', '.js', '.ts', '.java', '.c', '.cpp', '.php', '.rb', '.tex',
        '.go', '.rs', '.swift', '.kt', '.scala', '.r', '.m', '.pl'
    ],
    # 스크립트 파일
    'script': ['.sh', '.bash', '.zsh', '.fish', '.ps1', '.bat', '.cmd'],
    # 웹 개발
    'web': ['.html', '.css'],
    # 문서 및 마크업
    'markup': ['.md', '.txt', '.rst', '.adoc', '.org'],
    # 노트북 형식
    'notebook': ['.ipynb'],
    # 데이터 및 설정
    'config': [
        '.json', '.xml', '.yaml', '.yml', '.toml', '.ini', '.conf',
        '.config', '.env', '.properties'
    ],
    # 빌드 및 프로젝트
    'build': ['.makefile', '.gradle', '.cmake'],
    # 시스템 및 분석
    'system': ['.log', '.sql'],
}
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB

router = APIRouter()


# 지원 모델 목록
AVAILABLE_MODELS = [
    {"id": "gpt4o", "name": "GPT-4o", "provider": "OpenAI"},
    {"id": "gpt4om", "name": "GPT-4o-mini", "provider": "OpenAI"},
    {"id": "gpt41", "name": "GPT-4.1", "provider": "OpenAI"},
    {"id": "gpt5", "name": "GPT-5", "provider": "OpenAI"},
    {"id": "gpt52", "name": "GPT-5.2", "provider": "OpenAI"},
    {"id": "gpt5m", "name": "GPT-5 mini", "provider": "OpenAI"},
    {"id": "gpt5n", "name": "GPT-5 nano", "provider": "OpenAI"},
    {"id": "gpto3", "name": "O3", "provider": "OpenAI"},
    {"id": "gpto3p", "name": "O3-Pro", "provider": "OpenAI"},
    {"id": "gmn30", "name": "Gemini 3.0 Pro", "provider": "Google"},
    {"id": "gmn25f", "name": "Gemini 2.5 Flash", "provider": "Google"},
    {"id": "gmn25", "name": "Gemini 2.5 Pro", "provider": "Google"},
    {"id": "cld45o", "name": "Claude 4.5 Opus", "provider": "Anthropic"},
    {"id": "cld45s", "name": "Claude 4.5 Sonnet", "provider": "Anthropic"},
    {"id": "cld4o", "name": "Claude 4 Opus", "provider": "Anthropic"},
    # Perplexity 모델 임시 비활성화
    # {"id": "pplx", "name": "Sonar Reasoning", "provider": "Perplexity"},
    # {"id": "pplxp", "name": "Sonar Pro", "provider": "Perplexity"},
]


class ChatRequest(BaseModel):
    message: str
    model: str = "gpt4o"
    session_id: Optional[str] = None
    temperature: float = 0.7
    attachment_ids: Optional[List[int]] = None  # 첨부파일 ID 목록
    agent_mode: bool = False  # Agent 모드 활성화 여부 (실시간 검색, RAG 등)


class SessionRenameRequest(BaseModel):
    title: str


def get_file_type(filename: str) -> str:
    """파일 확장자로 타입 결정"""
    ext = os.path.splitext(filename)[1].lower()
    for file_type, extensions in ALLOWED_EXTENSIONS.items():
        if ext in extensions:
            return file_type
    return 'other'


def is_allowed_file(filename: str) -> bool:
    """허용된 파일인지 확인"""
    ext = os.path.splitext(filename)[1].lower()
    all_extensions = []
    for extensions in ALLOWED_EXTENSIONS.values():
        all_extensions.extend(extensions)
    return ext in all_extensions


def extract_file_content(file_path: str, original_filename: str) -> str:
    """파일에서 텍스트 내용 추출 (확장된 지원)"""
    ext = os.path.splitext(original_filename)[1].lower()
    content = ""

    # 텍스트 기반 파일 확장자 목록
    TEXT_EXTENSIONS = [
        # 문서 및 마크업
        '.txt', '.md', '.rst', '.adoc', '.org',
        # 표 및 데이터 (텍스트 형식)
        '.csv', '.tsv',
        # 프로그래밍 언어
        '.py', '.js', '.ts', '.java', '.c', '.cpp', '.php', '.rb', '.tex',
        '.go', '.rs', '.swift', '.kt', '.scala', '.r', '.m', '.pl',
        # 스크립트 파일
        '.sh', '.bash', '.zsh', '.fish', '.ps1', '.bat', '.cmd',
        # 웹 개발
        '.html', '.css',
        # 데이터 및 설정
        '.json', '.xml', '.yaml', '.yml', '.toml', '.ini', '.conf',
        '.config', '.env', '.properties',
        # 빌드 및 프로젝트
        '.makefile', '.gradle', '.cmake',
        # 시스템 및 분석
        '.log', '.sql',
    ]

    try:
        # PDF 파일
        if ext == '.pdf' and PDF_AVAILABLE:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                pages = []
                for i, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text:
                        pages.append(f"[페이지 {i+1}]\n{text}")
                content = "\n\n".join(pages)

        # Word 문서
        elif ext == '.docx' and DOCX_AVAILABLE:
            doc = DocxDocument(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = "\n".join(paragraphs)

        # Excel 파일
        elif ext in ['.xlsx', '.xls'] and XLSX_AVAILABLE:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            sheets_content = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_str = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_str.strip():
                        rows.append(row_str)
                if rows:
                    sheets_content.append(f"[시트: {sheet_name}]\n" + "\n".join(rows))
            content = "\n\n".join(sheets_content)

        # PowerPoint 파일
        elif ext in ['.pptx', '.ppt'] and PPTX_AVAILABLE:
            prs = Presentation(file_path)
            slides_content = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    slides_content.append(f"[슬라이드 {i}]\n" + "\n".join(slide_texts))
            content = "\n\n".join(slides_content)

        # HWP/HWPX 파일 (한글 문서)
        elif ext in ['.hwp', '.hwpx']:
            try:
                import olefile
                if ext == '.hwp' and olefile.isOleFile(file_path):
                    ole = olefile.OleFileIO(file_path)
                    if ole.exists('PrvText'):
                        prvtext = ole.openstream('PrvText').read()
                        content = prvtext.decode('utf-16', errors='ignore')
                    ole.close()
                else:
                    content = "[HWP/HWPX 파일은 텍스트 미리보기가 제한적입니다]"
            except ImportError:
                content = "[HWP 파일 지원을 위해 olefile 라이브러리가 필요합니다]"
            except Exception as e:
                content = f"[HWP 파일 읽기 오류: {str(e)}]"

        # Jupyter Notebook
        elif ext == '.ipynb':
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    notebook = json.load(f)
                cells_content = []
                for i, cell in enumerate(notebook.get('cells', []), 1):
                    cell_type = cell.get('cell_type', 'unknown')
                    source = ''.join(cell.get('source', []))
                    if source.strip():
                        cells_content.append(f"[셀 {i} - {cell_type}]\n{source}")
                content = "\n\n".join(cells_content)
            except Exception as e:
                content = f"[Jupyter Notebook 읽기 오류: {str(e)}]"

        # 텍스트 기반 파일 (확장된 목록)
        elif ext in TEXT_EXTENSIONS:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

        # 내용이 너무 길면 자르기 (약 50000자)
        if len(content) > 50000:
            content = content[:50000] + "\n\n... (내용이 너무 길어 일부만 표시됩니다)"

    except Exception as e:
        content = f"[파일 내용을 읽을 수 없습니다: {str(e)}]"

    return content


def is_image_file(filename: str) -> bool:
    """이미지 파일인지 확인"""
    ext = os.path.splitext(filename)[1].lower()
    return ext in ALLOWED_EXTENSIONS.get('image', [])


def get_image_base64(file_path: str) -> str:
    """이미지를 base64로 인코딩"""
    with open(file_path, 'rb') as f:
        return base64.b64encode(f.read()).decode('utf-8')


@router.get("/models")
async def get_models():
    """사용 가능한 LLM 모델 목록 반환"""
    return {"models": AVAILABLE_MODELS}


# ============================================================================
# 사용량 제한 API (일반 사용자용)
# ============================================================================

from utils.usage_limits import get_limits_summary, check_user_limits, check_model_access


@router.get("/usage")
async def get_my_usage(current_user: dict = Depends(get_current_user)):
    """
    현재 사용자의 사용량 및 제한 조회

    일반 사용자가 자신의 사용량과 남은 한도를 확인할 수 있습니다.
    """
    user_id = current_user["id"]

    summary = await get_limits_summary(user_id)

    return {
        "user_id": user_id,
        **summary
    }


@router.get("/usage/check")
async def check_my_limits(current_user: dict = Depends(get_current_user)):
    """
    현재 사용자가 메시지를 보낼 수 있는지 확인

    채팅 전에 호출하여 제한에 걸리는지 미리 확인할 수 있습니다.
    """
    user_id = current_user["id"]

    result = await check_user_limits(user_id)

    return result


@router.post("/send")
async def send_message(request: ChatRequest, current_user: dict = Depends(get_current_user)):
    """채팅 메시지 전송 (스트리밍 응답) - 인증 필요"""
    user_id = current_user["id"]

    # 세션 생성 또는 확인
    session_id = request.session_id
    if not session_id:
        # 새 세션 생성
        session_id = str(uuid.uuid4())
        await create_session(session_id, user_id, model=request.model)
    else:
        # 기존 세션 확인
        session = await get_session(session_id)
        if not session:
            raise HTTPException(status_code=404, detail="세션을 찾을 수 없습니다")
        if session["user_id"] != user_id:
            raise HTTPException(status_code=403, detail="접근 권한이 없습니다")

    # 첨부파일 내용 추출 (텍스트 + 이미지 분리)
    text_contents = []
    image_data_list = []  # (base64, mime_type, filename)

    if request.attachment_ids:
        for att_id in request.attachment_ids:
            att = await get_attachment(att_id)
            if att and att["user_id"] == user_id:
                if is_image_file(att["original_filename"]):
                    # 이미지는 base64로 인코딩
                    try:
                        img_base64 = get_image_base64(att["file_path"])
                        ext = os.path.splitext(att["original_filename"])[1].lower()
                        mime_map = {'.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', '.png': 'image/png',
                                   '.gif': 'image/gif', '.webp': 'image/webp', '.bmp': 'image/bmp'}
                        mime_type = mime_map.get(ext, 'image/png')
                        image_data_list.append((img_base64, mime_type, att["original_filename"]))
                    except Exception as e:
                        text_contents.append(f"[이미지 로드 실패: {att['original_filename']} - {str(e)}]")
                else:
                    # 문서 파일 내용 추출
                    content = extract_file_content(att["file_path"], att["original_filename"])
                    if content:
                        text_contents.append(
                            f"=== 첨부파일: {att['original_filename']} ===\n{content}\n=== 파일 끝 ==="
                        )

    # 텍스트 첨부파일 내용을 메시지에 추가
    full_text_message = request.message
    if text_contents:
        full_text_message = "\n\n".join(text_contents) + "\n\n---\n\n사용자 질문: " + request.message

    # 사용자 메시지 저장 (원본 메시지만) - ID 반환받기
    user_message_id = await add_message(session_id, "user", request.message)

    # 첫 메시지인 경우 세션 제목 설정
    messages = await get_session_messages(session_id)
    if len(messages) == 1:
        # 첫 질문의 앞 30자를 제목으로
        title = request.message[:30] + ("..." if len(request.message) > 30 else "")
        await update_session_title(session_id, title)

    # LLM 설정
    user_settings = {
        "model_provider": request.model,
        "temperature": request.temperature
    }

    # 비동기 제너레이터를 위한 변수 캡처
    _session_id = session_id
    _user_id = user_id
    _user_message_id = user_message_id
    _messages = messages
    _full_text_message = full_text_message
    _text_contents = text_contents
    _image_data_list = image_data_list
    _user_query = request.message
    _agent_mode = request.agent_mode
    _has_attachments = bool(request.attachment_ids)

    # ============ Agent 모드 처리 ============
    if _agent_mode:
        async def generate_agent():
            """Agent 모드 스트리밍 응답 생성"""
            try:
                # 첫 이벤트: 세션 ID와 사용자 메시지 ID 전송 (스트리밍 중단 시 정리용)
                yield f"data: {json.dumps({'session_id': _session_id, 'user_message_id': _user_message_id, 'done': False}, ensure_ascii=False)}\n\n"

                agent_service = get_agent_service()

                # 대화 히스토리 구성
                conversation_history = [
                    {"role": msg["role"], "content": msg["content"]}
                    for msg in _messages[:-1]  # 현재 메시지 제외
                ]

                full_response = ""
                tool_used = None

                # Agent 파이프라인 실행
                async for event in agent_service.process_with_agent(
                    user_input=_full_text_message,
                    tenant_id="default",
                    user_id=str(_user_id),
                    conversation_history=conversation_history,
                    user_settings=user_settings,
                    has_uploaded_docs=_has_attachments,
                    session_id=_session_id
                ):
                    event_type = event.get("type")

                    if event_type == "status":
                        # 상태 메시지 전송
                        yield f"data: {json.dumps({'status': event.get('message'), 'done': False, 'session_id': _session_id}, ensure_ascii=False)}\n\n"

                    elif event_type == "token":
                        # 응답 토큰 전송
                        content = event.get("content", "")
                        if content:
                            full_response += content
                            yield f"data: {json.dumps({'content': content, 'done': False, 'session_id': _session_id}, ensure_ascii=False)}\n\n"

                    elif event_type == "done":
                        tool_used = event.get("tool_used")

                    elif event_type == "error":
                        yield f"data: {json.dumps({'error': event.get('message'), 'done': True}, ensure_ascii=False)}\n\n"
                        return

                # 어시스턴트 응답 DB 저장
                message_id = await add_message(_session_id, "assistant", full_response)

                # 메모리 저장 (백그라운드)
                try:
                    memory_svc = await get_memory_service()
                    if memory_svc and memory_svc.is_available():
                        conversation_text = f"User: {_user_query}\nAssistant: {full_response[:500]}"
                        asyncio.create_task(
                            memory_svc.store_memory(
                                tenant_id="default",
                                user_id=_user_id,
                                text=conversation_text,
                                category="history"
                            )
                        )
                except Exception as e:
                    logger.warning(f"Agent memory storage failed: {e}")

                # 완료 이벤트
                yield f"data: {json.dumps({'content': '', 'done': True, 'session_id': _session_id, 'message_id': message_id, 'tool_used': tool_used}, ensure_ascii=False)}\n\n"

            except Exception as e:
                logger.error(f"Agent mode error: {e}")
                yield f"data: {json.dumps({'error': str(e), 'done': True}, ensure_ascii=False)}\n\n"

        return StreamingResponse(
            generate_agent(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "X-Accel-Buffering": "no"
            }
        )

    # ============ 일반 모드 (개인화 향상 적용) ============
    async def generate():
        try:
            llm = get_llm(user_settings)

            # ============ 개인화 향상 모듈 초기화 ============
            enhancer = None
            query_intent = None
            skip_rag = False

            try:
                enhancer = await get_personalization_enhancer()
            except Exception as e:
                logger.warning(f"Personalization enhancer init failed: {e}")

            # ============ 1. 시맨틱 라우터 - 쿼리 의도 분류 ============
            if enhancer and SEMANTIC_ROUTER_ENABLED:
                try:
                    query_intent = await enhancer.classify_query(_user_query)
                    logger.info(f"[{"default"}] User {_user_id}: Query intent = {query_intent.value}")

                    # 인사/잡담은 RAG 스킵
                    if query_intent == QueryIntent.GREETING:
                        skip_rag = True
                        logger.info(f"[{"default"}] Skipping RAG for greeting query")
                except Exception as e:
                    logger.warning(f"Semantic router failed: {e}")

            # ============ 2. 개인화 컨텍스트 조회 (Qdrant) ============
            personalized_context = ""
            relevance_ok = True
            memory_svc = None

            try:
                memory_svc = await get_memory_service()
                logger.info(f"[{"default"}] Memory service initialized: available={memory_svc.is_available() if memory_svc else 'None'}, skip_rag={skip_rag}")
                if memory_svc and memory_svc.is_available() and not skip_rag:
                    personalized_context = await memory_svc.get_personalized_context(
                        tenant_id="default",
                        user_id=_user_id,
                        current_query=_user_query
                    )
                    if personalized_context:
                        logger.info(f"[{"default"}] User {_user_id}: Personalized context loaded ({len(personalized_context)} chars)")
                    else:
                        logger.info(f"[{"default"}] User {_user_id}: No personalized context found")

                    # ============ 3. Self-RAG - 관련성 검증 ============
                    if enhancer and SELF_RAG_ENABLED and personalized_context:
                        try:
                            relevance_result = await enhancer.check_relevance(
                                query=_user_query,
                                context=personalized_context
                            )
                            relevance_ok = relevance_result.is_relevant

                            if not relevance_ok:
                                logger.info(f"[{"default"}] Self-RAG: Context not relevant (score={relevance_result.score:.2f}), using fallback")
                                # 관련성 낮으면 컨텍스트 축소 또는 제거
                                if relevance_result.score < 0.3:
                                    personalized_context = ""  # 너무 낮으면 제거
                                else:
                                    # 부분적으로 관련 있으면 간략화
                                    personalized_context = personalized_context[:500] + "..."
                        except Exception as e:
                            logger.warning(f"Self-RAG check failed: {e}")
            except Exception as e:
                logger.warning(f"Personalization context retrieval failed: {e}")

            # ============ MH_RAG 문서 검색 ============
            mh_rag_context = ""
            if MH_RAG_ENABLED and not skip_rag:
                try:
                    from services.mh_rag_service import needs_mh_rag_search
                    if needs_mh_rag_search(_user_query):
                        mh_rag_svc = await get_mh_rag_service()
                        if mh_rag_svc and mh_rag_svc.is_available():
                            documents = await mh_rag_svc.search(
                                query=_user_query,
                                top_k=5
                            )
                            if documents:
                                mh_rag_context = mh_rag_svc.format_context(documents)
                                logger.info(f"[{"default"}] MH_RAG: {len(documents)}개 문서 검색 완료")
                except Exception as e:
                    logger.warning(f"MH_RAG search failed: {e}")

            # ============ 6. 계층적 메모리 컨텍스트 구성 ============
            hierarchical_context = ""
            if enhancer and HIERARCHICAL_MEMORY_ENABLED:
                try:
                    h_memory = await enhancer.get_hierarchical_memory(
                        tenant_id="default",
                        user_id=str(_user_id)
                    )
                    if h_memory:
                        context_parts = []
                        # 장기 기억 (사용자 프로필)
                        if h_memory.long_term and h_memory.long_term.interests:
                            context_parts.append(f"[사용자 관심사] {', '.join(h_memory.long_term.interests[:5])}")
                        if h_memory.long_term and h_memory.long_term.preferences:
                            prefs = [f"{k}:{v}" for k, v in list(h_memory.long_term.preferences.items())[:3]]
                            context_parts.append(f"[선호도] {', '.join(prefs)}")
                        # 중기 기억 (세션 요약)
                        if h_memory.mid_term:
                            recent_summaries = h_memory.mid_term[-2:]  # 최근 2개 세션
                            for ss in recent_summaries:
                                if ss.summary:
                                    context_parts.append(f"[이전 대화] {ss.summary[:100]}")
                        if context_parts:
                            hierarchical_context = "\n".join(context_parts)
                            logger.info(f"[{"default"}] User {_user_id}: Hierarchical context loaded")
                except Exception as e:
                    logger.warning(f"Hierarchical memory retrieval failed: {e}")

            # ============ 7. Long-term 프로파일 로드 (DB 기반) ============
            longterm_profile_context = ""
            if LONGTERM_MEMORY_ENABLED:
                try:
                    profile_svc = await get_profile_service()
                    profile_summary = await profile_svc.get_profile_summary_for_llm(
                        tenant_id="default",
                        user_id=_user_id
                    )
                    if profile_summary:
                        longterm_profile_context = profile_summary
                        logger.info(f"[{"default"}] User {_user_id}: Long-term profile loaded")
                except Exception as e:
                    logger.warning(f"Long-term profile retrieval failed: {e}")

            # 대화 이력을 메시지 형식으로 변환
            chat_messages = []

            # 시스템 프롬프트 (마크다운 형식 지시 + 간결함 유도)
            base_system_prompt = """당신은 한림대학교 구성원을 위한 AI 어시스턴트입니다.

## 핵심 원칙
- 질문에 직접적이고 명확하게 답변
- 불필요한 아첨/칭찬 금지 (예: "좋은 질문입니다!")
- "더 궁금한 점 있으시면...", "~해드릴까요?" 등 헤징 표현 자제

## 응답 길이 가이드
- 단순 질문(정의, 개념): 2-4문장으로 간결하게
- 기술적 질문(코드, 절차): 핵심만 설명, 예시 1개
- 복잡한 질문: 구조화하되 부연 설명 자제

## 코드 포맷팅 규칙 (필수!)
1. 인라인 코드(변수명, 키워드): 백틱 사용 - `for`, `while`, `print()`
2. 코드 예제(2줄 이상): 반드시 코드블록(```) 사용하고 줄바꿈+들여쓰기 포함
   - 올바른 예:
   ```python
   if x > 0:
       print("양수")
   ```
   - 잘못된 예: if x > 0: print("양수") (한 줄로 작성 금지)
3. 코드 예제는 절대 한 줄로 쓰지 말 것"""

            # 개인화 컨텍스트 통합 (Long-term 프로파일 + 계층적 메모리 + 개인화 + MH_RAG)
            combined_context = ""
            if longterm_profile_context:
                combined_context += f"{longterm_profile_context}\n\n"
            if hierarchical_context:
                combined_context += f"=== 사용자 기억 ===\n{hierarchical_context}\n\n"
            if personalized_context:
                combined_context += f"=== 개인화 정보 ===\n{personalized_context}\n\n"
            if mh_rag_context:
                combined_context += f"=== 참고 문서 ===\n{mh_rag_context}"

            if combined_context.strip():
                system_prompt = f"""{base_system_prompt}

아래는 질문에 답변하기 위한 참고 자료입니다. 이 정보를 바탕으로 정확하게 답변하세요:

{combined_context.strip()}

위 정보를 활용하여 사용자 질문에 답변해 주세요. 참고 자료에 없는 내용은 일반 지식으로 보완하되, 출처가 있는 정보는 정확히 인용하세요."""
                chat_messages.append(SystemMessage(content=system_prompt))
            else:
                chat_messages.append(SystemMessage(content=base_system_prompt))
            for i, msg in enumerate(_messages):
                if msg["role"] == "user":
                    # 마지막 메시지(현재 질문)에만 첨부파일 포함
                    if i == len(_messages) - 1 and (_text_contents or _image_data_list):
                        # 이미지가 있으면 멀티모달 메시지 생성
                        if _image_data_list:
                            content_parts = []
                            # 텍스트 먼저 추가
                            content_parts.append({"type": "text", "text": _full_text_message})
                            # 이미지들 추가
                            for img_b64, mime_type, filename in _image_data_list:
                                content_parts.append({
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:{mime_type};base64,{img_b64}"
                                    }
                                })
                            chat_messages.append(HumanMessage(content=content_parts))
                        else:
                            # 텍스트만 있는 경우
                            chat_messages.append(HumanMessage(content=_full_text_message))
                    else:
                        chat_messages.append(HumanMessage(content=msg["content"]))
                else:
                    chat_messages.append(AIMessage(content=msg["content"]))

            full_response = ""

            # 스트리밍 응답
            async for chunk in llm.astream(chat_messages):
                content = chunk.content if hasattr(chunk, 'content') else str(chunk)
                if content:
                    full_response += content
                    yield f"data: {json.dumps({'content': content, 'done': False, 'session_id': _session_id}, ensure_ascii=False)}\n\n"

            # 어시스턴트 응답 DB 저장 (message_id 반환)
            message_id = await add_message(_session_id, "assistant", full_response)

            # ============ 개인화 메모리 저장 (비동기 백그라운드) ============
            try:
                if memory_svc and memory_svc.is_available():
                    # 대화 히스토리 저장
                    conversation_text = f"User: {_user_query}\nAssistant: {full_response[:500]}"
                    asyncio.create_task(
                        memory_svc.store_memory(
                            tenant_id="default",
                            user_id=_user_id,
                            text=conversation_text,
                            category="history"
                        )
                    )

                    # 사용자 정보 추출 및 저장 (백그라운드)
                    asyncio.create_task(
                        memory_svc.extract_and_store_user_info(
                            tenant_id="default",
                            user_id=_user_id,
                            user_message=_user_query,
                            assistant_response=full_response
                        )
                    )
                    logger.debug(f"[{"default"}] User {_user_id}: Memory storage tasks scheduled")
            except Exception as e:
                logger.warning(f"Memory storage scheduling failed: {e}")

            # ============ 4. LLM-as-Judge 평가 (백그라운드) ============
            if enhancer and LLM_JUDGE_ENABLED:
                async def run_judge_evaluation():
                    try:
                        judge_score = await enhancer.evaluate_response(
                            query=_user_query,
                            response=full_response,
                            context=personalized_context if personalized_context else None
                        )
                        logger.info(
                            f"[{"default"}] LLM-Judge scores - "
                            f"Accuracy: {judge_score.accuracy}/5, "
                            f"Helpfulness: {judge_score.helpfulness}/5, "
                            f"Personalization: {judge_score.personalization}/5, "
                            f"Friendliness: {judge_score.friendliness}/5, "
                            f"Overall: {judge_score.overall:.2f}/5"
                        )

                        # 평가 점수를 메모리에 저장 (향후 분석용)
                        if memory_svc and memory_svc.is_available():
                            eval_text = (
                                f"[평가 결과] 정확성:{judge_score.accuracy}, "
                                f"도움됨:{judge_score.helpfulness}, "
                                f"개인화:{judge_score.personalization}, "
                                f"친절도:{judge_score.friendliness}"
                            )
                            await memory_svc.store_memory(
                                tenant_id="default",
                                user_id=_user_id,
                                text=eval_text,
                                category="evaluation"
                            )
                    except Exception as e:
                        logger.warning(f"LLM-Judge evaluation failed: {e}")

                asyncio.create_task(run_judge_evaluation())

            # ============ 5. 계층적 메모리 업데이트 (백그라운드) ============
            if enhancer and HIERARCHICAL_MEMORY_ENABLED:
                async def update_hierarchical_memory():
                    try:
                        await enhancer.update_memory(
                            tenant_id="default",
                            user_id=str(_user_id),
                            query=_user_query,
                            response=full_response,
                            query_intent=query_intent
                        )
                        logger.debug(f"[{"default"}] User {_user_id}: Hierarchical memory updated")
                    except Exception as e:
                        logger.warning(f"Hierarchical memory update failed: {e}")

                asyncio.create_task(update_hierarchical_memory())

            # ============ 백그라운드: Long-term 프로파일 자동 추출 ============
            if PROFILE_AUTO_EXTRACT and len(_messages) >= 3:  # 최소 3턴 이상일 때만
                async def extract_longterm_profile():
                    try:
                        profile_svc = await get_profile_service()
                        result = await profile_svc.extract_profile_from_conversation(
                            tenant_id="default",
                            user_id=_user_id,
                            session_id=_session_id,
                            messages=_messages
                        )
                        if result.get("extracted") and result.get("facts_count", 0) > 0:
                            logger.info(f"[{"default"}] User {_user_id}: Extracted {result['facts_count']} facts to Long-term memory")
                    except Exception as e:
                        logger.warning(f"Long-term profile extraction failed: {e}")

                asyncio.create_task(extract_longterm_profile())

            yield f"data: {json.dumps({'content': '', 'done': True, 'session_id': _session_id, 'message_id': message_id}, ensure_ascii=False)}\n\n"

        except Exception as e:
            error_msg = f"Error: {str(e)}"
            yield f"data: {json.dumps({'error': error_msg, 'done': True}, ensure_ascii=False)}\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no"
        }
    )


@router.get("/sessions")
async def list_sessions(current_user: dict = Depends(get_current_user)):
    """사용자의 대화 세션 목록 조회"""
    sessions = await get_user_sessions(current_user["id"])
    return {"sessions": sessions}


@router.get("/sessions/{session_id}")
async def get_session_detail(session_id: str, current_user: dict = Depends(get_current_user)):
    """세션 상세 조회 (메시지 포함)"""
    session = await get_session(session_id)

    if not session:
        raise HTTPException(status_code=404, detail="세션을 찾을 수 없습니다")
    if session["user_id"] != current_user["id"]:
        raise HTTPException(status_code=403, detail="접근 권한이 없습니다")

    messages = await get_session_messages(session_id)

    return {
        "session": session,
        "messages": messages
    }


@router.put("/sessions/{session_id}")
async def rename_session(
    session_id: str,
    request: SessionRenameRequest,
    current_user: dict = Depends(get_current_user)
):
    """세션 이름 변경"""
    session = await get_session(session_id)

    if not session:
        raise HTTPException(status_code=404, detail="세션을 찾을 수 없습니다")
    if session["user_id"] != current_user["id"]:
        raise HTTPException(status_code=403, detail="접근 권한이 없습니다")

    await update_session_title(session_id, request.title)
    return {"status": "updated"}


@router.delete("/sessions/{session_id}")
async def delete_session_endpoint(session_id: str, current_user: dict = Depends(get_current_user)):
    """세션 삭제"""
    session = await get_session(session_id)

    if not session:
        raise HTTPException(status_code=404, detail="세션을 찾을 수 없습니다")
    if session["user_id"] != current_user["id"]:
        raise HTTPException(status_code=403, detail="접근 권한이 없습니다")

    await delete_session(session_id)
    return {"status": "deleted"}


# ============ 비인증 API (하위 호환용) ============

@router.post("/send/guest")
async def send_message_guest(request: ChatRequest):
    """게스트 채팅 (인증 불필요, 저장 안됨)"""
    user_settings = {
        "model_provider": request.model,
        "temperature": request.temperature
    }

    async def generate():
        try:
            llm = get_llm(user_settings)

            # 단일 메시지만 처리 (이력 없음)
            messages = [("human", request.message)]
            full_response = ""

            async for chunk in llm.astream(messages):
                content = chunk.content if hasattr(chunk, 'content') else str(chunk)
                if content:
                    full_response += content
                    yield f"data: {json.dumps({'content': content, 'done': False}, ensure_ascii=False)}\n\n"

            yield f"data: {json.dumps({'content': '', 'done': True}, ensure_ascii=False)}\n\n"

        except Exception as e:
            error_msg = f"Error: {str(e)}"
            yield f"data: {json.dumps({'error': error_msg, 'done': True}, ensure_ascii=False)}\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no"
        }
    )


# ============ 파일 업로드/다운로드 API ============

@router.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    session_id: Optional[str] = Form(None),
    current_user: dict = Depends(get_current_user)
):
    """파일 업로드 (이미지, 문서 등) - 스트리밍 방식으로 메모리 효율적"""
    import aiofiles

    user_id = current_user["id"]

    # 파일 검증
    if not file.filename:
        raise HTTPException(status_code=400, detail="파일명이 없습니다")

    if not is_allowed_file(file.filename):
        raise HTTPException(status_code=400, detail="허용되지 않는 파일 형식입니다")

    # 세션이 없으면 임시 세션 ID 생성
    if not session_id:
        session_id = f"temp_{str(uuid.uuid4())}"

    # 사용자별 디렉토리 생성
    user_upload_dir = os.path.join(UPLOAD_DIR, str(user_id))
    os.makedirs(user_upload_dir, exist_ok=True)

    # 고유 파일명 생성
    file_ext = os.path.splitext(file.filename)[1]
    unique_filename = f"{uuid.uuid4()}{file_ext}"
    file_path = os.path.join(user_upload_dir, unique_filename)

    # 스트리밍 방식으로 파일 저장 (메모리 효율적)
    file_size = 0
    chunk_size = 1024 * 1024  # 1MB 청크

    try:
        async with aiofiles.open(file_path, 'wb') as out_file:
            while True:
                chunk = await file.read(chunk_size)
                if not chunk:
                    break
                file_size += len(chunk)
                # 파일 크기 체크 (스트리밍 중)
                if file_size > MAX_FILE_SIZE:
                    await out_file.close()
                    os.remove(file_path)
                    raise HTTPException(
                        status_code=400,
                        detail=f"파일 크기가 {MAX_FILE_SIZE // (1024*1024)}MB를 초과합니다"
                    )
                await out_file.write(chunk)
    except HTTPException:
        raise
    except Exception as e:
        if os.path.exists(file_path):
            os.remove(file_path)
        raise HTTPException(status_code=500, detail=f"파일 저장 실패: {str(e)}")

    # DB에 기록
    file_type = get_file_type(file.filename)
    attachment_id = await add_attachment(
        session_id=session_id,
        user_id=user_id,
        filename=unique_filename,
        original_filename=file.filename,
        file_type=file_type,
        file_size=file_size,
        file_path=file_path
    )

    return {
        "id": attachment_id,
        "filename": unique_filename,
        "original_filename": file.filename,
        "file_type": file_type,
        "file_size": file_size,
        "session_id": session_id
    }


@router.get("/attachments/{attachment_id}")
async def get_attachment_info(attachment_id: int, current_user: dict = Depends(get_current_user)):
    """첨부파일 정보 조회"""
    attachment = await get_attachment(attachment_id)

    if not attachment:
        raise HTTPException(status_code=404, detail="첨부파일을 찾을 수 없습니다")
    if attachment["user_id"] != current_user["id"]:
        raise HTTPException(status_code=403, detail="접근 권한이 없습니다")

    return attachment


@router.get("/attachments/{attachment_id}/download")
async def download_attachment(attachment_id: int, current_user: dict = Depends(get_current_user)):
    """첨부파일 다운로드"""
    attachment = await get_attachment(attachment_id)

    if not attachment:
        raise HTTPException(status_code=404, detail="첨부파일을 찾을 수 없습니다")
    if attachment["user_id"] != current_user["id"]:
        raise HTTPException(status_code=403, detail="접근 권한이 없습니다")

    file_path = attachment["file_path"]
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="파일이 존재하지 않습니다")

    return FileResponse(
        path=file_path,
        filename=attachment["original_filename"],
        media_type="application/octet-stream"
    )


@router.get("/sessions/{session_id}/attachments")
async def list_session_attachments(session_id: str, current_user: dict = Depends(get_current_user)):
    """세션의 모든 첨부파일 목록"""
    session = await get_session(session_id)

    if session and session["user_id"] != current_user["id"]:
        raise HTTPException(status_code=403, detail="접근 권한이 없습니다")

    attachments = await get_session_attachments(session_id)
    return {"attachments": attachments}


class FeedbackRequest(BaseModel):
    message_id: int
    feedback: str  # 'like' or 'dislike' or null


class PartialMessageRequest(BaseModel):
    session_id: str
    content: str
    model: Optional[str] = None


@router.post("/messages/partial")
async def save_partial_message(request: PartialMessageRequest, current_user: dict = Depends(get_current_user)):
    """
    스트리밍 중단 시 부분 응답 저장

    사용자가 세션을 전환하거나 스트리밍을 중단했을 때
    현재까지 받은 응답을 저장합니다.
    """
    if not request.content or not request.content.strip():
        return {"success": False, "message": "빈 내용은 저장하지 않습니다"}

    try:
        # 세션 소유권 확인
        async with get_db() as session:
            result = await session.execute(
                text("SELECT user_id FROM sessions WHERE id = :session_id"),
                {"session_id": request.session_id}
            )
            sess = result.fetchone()

            if not sess:
                return {"success": False, "message": "세션을 찾을 수 없습니다"}

            if sess[0] != current_user["id"]:
                raise HTTPException(status_code=403, detail="접근 권한이 없습니다")

        # 부분 응답 저장
        message_id = await add_message(request.session_id, "assistant", request.content.strip())

        logger.info(f"Partial response saved for session {request.session_id}: {len(request.content)} chars")

        return {"success": True, "message_id": message_id}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to save partial message: {e}")
        return {"success": False, "message": str(e)}


@router.delete("/messages/{message_id}")
async def delete_message(message_id: int, current_user: dict = Depends(get_current_user)):
    """
    특정 메시지 삭제

    스트리밍이 중단되어 응답이 없는 경우,
    해당 사용자 메시지를 삭제하여 일관성 유지
    """
    try:
        async with get_db() as session:
            # 메시지 존재 및 소유권 확인
            result = await session.execute(
                text("""
                    SELECT m.id, m.session_id, s.user_id
                    FROM messages m
                    JOIN sessions s ON m.session_id = s.id
                    WHERE m.id = :message_id
                """),
                {"message_id": message_id}
            )
            msg = result.fetchone()

            if not msg:
                return {"success": False, "message": "메시지를 찾을 수 없습니다"}

            if msg[2] != current_user["id"]:
                raise HTTPException(status_code=403, detail="접근 권한이 없습니다")

            # 메시지 삭제
            await session.execute(
                text("DELETE FROM messages WHERE id = :msg_id"),
                {"msg_id": message_id}
            )

            logger.info(f"Deleted message {message_id} from session {msg[1]}")

            return {"success": True, "deleted_message_id": message_id}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to delete message: {e}")
        return {"success": False, "message": str(e)}


@router.post("/messages/feedback")
async def update_message_feedback(request: FeedbackRequest, current_user: dict = Depends(get_current_user)):
    """메시지 피드백 저장"""
    # 피드백 값 검증
    if request.feedback not in ['like', 'dislike', None, '']:
        raise HTTPException(status_code=400, detail="피드백은 'like', 'dislike' 또는 빈 값이어야 합니다")

    feedback_value = request.feedback if request.feedback else None

    try:
        async with get_db() as session:
            # 메시지 존재 여부 확인
            result = await session.execute(
                text("SELECT id, session_id FROM messages WHERE id = :message_id"),
                {"message_id": request.message_id}
            )
            message = result.fetchone()

            if not message:
                raise HTTPException(status_code=404, detail="메시지를 찾을 수 없습니다")

            # 해당 세션이 현재 사용자의 것인지 확인
            result = await session.execute(
                text("SELECT user_id FROM sessions WHERE id = :session_id"),
                {"session_id": message[1]}
            )
            sess = result.fetchone()

            if not sess or sess[0] != current_user["id"]:
                raise HTTPException(status_code=403, detail="접근 권한이 없습니다")

            # 피드백 업데이트
            await session.execute(
                text("UPDATE messages SET feedback = :feedback WHERE id = :message_id"),
                {"feedback": feedback_value, "message_id": request.message_id}
            )

            logger.info(f"User {current_user['id']}: Feedback '{feedback_value}' for message {request.message_id}")

            return {"success": True, "message_id": request.message_id, "feedback": feedback_value}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Feedback update error: {e}")
        raise HTTPException(status_code=500, detail="피드백 저장 실패")
